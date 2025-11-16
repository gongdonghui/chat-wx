from app.api import *
import requests
from tasks import generate_meeting_minutes

from rank_bm25 import BM25Okapi
import faiss
import numpy as np
import numpy as np
import re
import requests
import json
import configparser
import os
import tempfile
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
from werkzeug.utils import secure_filename
import uuid
from app.utils.storage import storage

rag = Blueprint('rag', __name__)

# 加载配置文件
config = configparser.ConfigParser()
config.read('config.ini')

# 加载配置
VLLM_SERVER_URL = config.get('vllm', 'server_url', fallback='http://your-vllm-server-address/generate')
VLLM_RERANK_URL = config.get('vllm', 'rerank_url', fallback='http://your-vllm-server-address/rerank')
EMBEDDING_MODEL = 'BAAI/bge-large-zh'  # 更适合中文的嵌入模型
CHUNK_SIZE = 512
CHUNK_OVERLAP = 50
RRF_K = 60
TOP_K = 10
TOP_N = 5
RERANK_TOP_N = 5  # rerank后保留的文档数

# VLLM嵌入服务配置
VLLM_EMBEDDING_URL = "http://localhost:8000/embed"  # VLLM嵌入服务地址

documents = []
embeddings = []
bm25 = None
index = None

# 文本分块函数
def split_text(text, chunk_size=CHUNK_SIZE, chunk_overlap=CHUNK_OVERLAP, mode='fixed'):
    """
    将长文本分割成块
    :param text: 输入文本
    :param chunk_size: 块大小（仅在fixed模式下有效）
    :param chunk_overlap: 滑动窗口重叠大小
    :param mode: 分割模式，可选值：'fixed'（固定长度）, 'paragraph'（按段分割）
    :return: 分块后的文本列表
    """
    if mode == 'paragraph':
        # 按段分割，保留段落结构
        chunks = re.split(r'(\n\n+)', text.strip())
        # 合并段落分隔符到前一个块
        merged_chunks = []
        for i in range(0, len(chunks), 2):
            chunk = chunks[i].strip()
            if i + 1 < len(chunks):
                chunk += chunks[i+1]  # 保留段落分隔符
            if chunk:
                merged_chunks.append(chunk)
        return merged_chunks
    
    # 默认固定长度分割，带滑动窗口，保留标点符号
    # 使用正则分割并保留分隔符
    sentences = re.split(r'([。！？\n])', text)
    chunks = []
    current_chunk = []
    current_length = 0
    
    # 将分割后的句子和标点符号重新组合
    combined_sentences = []
    for i in range(0, len(sentences), 2):
        if i + 1 < len(sentences):
            combined_sentences.append(sentences[i] + sentences[i+1])
        else:
            combined_sentences.append(sentences[i])
    
    for sentence in combined_sentences:
        if not sentence:
            continue
        sentence_length = len(sentence)
        
        if current_length + sentence_length > chunk_size and current_chunk:
            chunks.append(''.join(current_chunk))
            
            # 实现滑动窗口：保留最后chunk_overlap个字符
            window_text = ''.join(current_chunk)
            if len(window_text) <= chunk_overlap:
                # 如果当前块长度小于等于重叠大小，则保留整个块
                current_chunk = [s for s in current_chunk if s]
            else:
                # 保留最后chunk_overlap个字符作为重叠
                overlap_part = window_text[-chunk_overlap:]
                # 确保重叠部分是完整的句子
                new_current_chunk = []
                temp_text = ''
                for sent in reversed(current_chunk):
                    temp_text = sent + temp_text
                    if temp_text in overlap_part or not new_current_chunk:  # 至少保留一个句子
                        new_current_chunk.insert(0, sent)
                    else:
                        break
                current_chunk = new_current_chunk if new_current_chunk else []
            
            current_length = sum(len(s) for s in current_chunk)
        
        current_chunk.append(sentence)
        current_length += sentence_length
    
    if current_chunk:
        chunks.append(''.join(current_chunk))
    
    # 移除空块
    chunks = [chunk.strip() for chunk in chunks if chunk.strip()]
    return chunks

# RRF融合排序
def rrf_fusion(results_list, k=60):
    """
    逆序位融合（支持多个检索结果）
    :param results_list: 检索结果列表 [[(doc_id, score), ...], ...]
    :param k: RRF参数
    :return: 融合后的结果列表 [(doc_id, score), ...]
    """
    rrf_scores = {}
    
    for results in results_list:
        for rank, (doc_id, score) in enumerate(results):
            if doc_id not in rrf_scores:
                rrf_scores[doc_id] = 0
            rrf_scores[doc_id] += 1 / (rank + k)
    
    # 按得分降序排序
    fused_results = sorted(rrf_scores.items(), key=lambda x: x[1], reverse=True)
    return fused_results

@rag.route('/rag/upload', methods=['POST'])
def upload_text():
    """
    上传文本并进行分块、索引
    """
    global documents, embeddings, bm25, index
    
    text = request.form.get('text')
    if not text:
        return jsonify({'error': 'No text provided'}), 400
    
    # 获取分块参数
    mode = request.form.get('mode', 'fixed')  # 默认为固定长度模式
    chunk_size = int(request.form.get('chunk_size', CHUNK_SIZE))  # 块大小
    chunk_overlap = int(request.form.get('chunk_overlap', CHUNK_OVERLAP))  # 滑动窗口重叠大小
    
    # 验证模式
    if mode not in ['fixed', 'paragraph']:
        return jsonify({'error': 'Invalid mode. Supported modes: fixed, paragraph'}), 400
    
    # 文本分块
    chunks = split_text(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap, mode=mode)
    if not chunks:
        return jsonify({'error': 'Text is too short or empty'}), 400
    
    # 更新文档列表
    documents.extend(chunks)
    
    # 生成文档向量（使用VLLM嵌入服务）
    headers = {"Content-Type": "application/json"}
    data = {"model": EMBEDDING_MODEL, "input": chunks}
    response = requests.post(VLLM_EMBEDDING_URL, json=data, headers=headers)
    response.raise_for_status()
    new_embeddings = np.array([item["embedding"] for item in response.json()["data"]], dtype='float32')
    embeddings.extend(new_embeddings.tolist())
    
    # 构建BM25索引（优化参数）
    tokenized_docs = [doc.split() for doc in documents]
    bm25 = BM25Okapi(tokenized_docs, k1=1.5, b=0.75)  # 调整BM25参数提高精度
    
    # 构建FAISS向量索引
    embeddings_np = np.array(embeddings, dtype='float32')
    dimension = embeddings_np.shape[1]
    index = faiss.IndexFlatL2(dimension)
    index.add(embeddings_np)
    
    return jsonify({
        'message': 'Text uploaded and indexed successfully',
        'num_chunks': len(chunks),
        'total_docs': len(documents)
    })

@rag.route('/rag/query', methods=['POST'])
def rag_query():
    """
    执行RAG查询，融合向量检索和BM25
    """
    global documents, embeddings, bm25, index
    
    query = request.form.get('query')
    if not query:
        return jsonify({'error': 'No query provided'}), 400
    
    if not documents or not bm25 or not index:
        return jsonify({'error': 'No documents indexed yet'}), 400
    
    # 参数配置
    TOP_N = 10  # 每个检索方法取前N个结果
    FINAL_TOP_N = 5  # 最终返回的相关文档数
    SIMILARITY_THRESHOLD = 0.2  # 向量检索相似度阈值（归一化后）
    
    # 1. 查询理解与扩展
    import sys
    import os
    sys.path.append('/Users/gongshuai/workspace/chat-wx')
    from query_expansion import expand_query
    expanded_query = expand_query(query)
    
    # 2. BM25检索
    tokenized_query = expanded_query.split()
    bm25_scores = bm25.get_scores(tokenized_query)
    bm25_results = [(i, score) for i, score in enumerate(bm25_scores)]
    bm25_results = sorted(bm25_results, key=lambda x: x[1], reverse=True)[:TOP_N]
    
    # 2. 向量检索
    headers = {"Content-Type": "application/json"}
    data = {"model": EMBEDDING_MODEL, "input": [query]}
    response = requests.post(VLLM_EMBEDDING_URL, json=data, headers=headers)
    response.raise_for_status()
    query_embedding = np.array([response.json()["data"][0]["embedding"]], dtype='float32')
    D, I = index.search(query_embedding, TOP_N)
    
    # 归一化向量检索得分（转换为相似度）
    vector_results = []
    for i, d in zip(I[0], D[0]):
        # 计算余弦相似度（假设FAISS使用的是L2距离）
        # 余弦相似度 = 1 / (1 + L2距离) （简化版，更精确的需要向量归一化）
        similarity = 1 / (1 + d)
        if similarity >= SIMILARITY_THRESHOLD:
            vector_results.append((i, similarity))
    
    # 3. 关键词匹配增强
    # 提取查询中的关键词
    import jieba
    jieba.setLogLevel(20)
    keywords = list(jieba.cut_for_search(query))
    keywords = [kw for kw in keywords if len(kw) > 1]  # 过滤短词
    
    keyword_results = []
    if keywords:
        for i, doc in enumerate(documents):
            score = 0
            for kw in keywords:
                if kw in doc:
                    score += 1
            if score > 0:
                keyword_results.append((i, score))
        keyword_results = sorted(keyword_results, key=lambda x: x[1], reverse=True)[:TOP_N]
    
    # 4. 融合排序
    results_list = [bm25_results, vector_results]
    if keyword_results:
        results_list.append(keyword_results)
    
    fused_results = rrf_fusion(results_list, k=50)  # 调整k值可能影响结果
    
    # 5. 去重并过滤低得分结果
    unique_docs = {}
    for doc_id, score in fused_results:
        if doc_id not in unique_docs:
            unique_docs[doc_id] = score
    
    # 6. Rerank重排
    reranked_docs = []
    candidate_docs = sorted(unique_docs.items(), key=lambda x: x[1], reverse=True)[:TOP_K]  # 取前TOP_K个结果进行rerank
    
    if candidate_docs:
        try:
            # 调用vllm的rerank服务
            rerank_input = {
                'query': query,
                'documents': [documents[doc_id] for doc_id, score in candidate_docs],
                'top_n': RERANK_TOP_N
            }
            
            response = requests.post(
                VLLM_RERANK_URL,
                json=rerank_input
            )
            
            if response.status_code == 200:
                rerank_results = response.json().get('results', [])
                
                # 构建rerank后的结果
                reranked_docs = []
                for i, result in enumerate(rerank_results):
                    doc_id = candidate_docs[i][0]
                    reranked_docs.append({
                        'doc_id': doc_id,
                        'content': documents[doc_id],
                        'score': result.get('score', candidate_docs[i][1]),
                        'rerank_score': result.get('score', 0.0)
                    })
                
                # 如果rerank返回结果不足，补充原始结果
                if len(reranked_docs) < FINAL_TOP_N:
                    remaining_docs = [doc for doc in candidate_docs if doc[0] not in [d['doc_id'] for d in reranked_docs]]
                    for doc_id, score in remaining_docs[:FINAL_TOP_N - len(reranked_docs)]:
                        reranked_docs.append({
                            'doc_id': doc_id,
                            'content': documents[doc_id],
                            'score': score
                        })
        except Exception as e:
            # 若rerank失败，使用原始融合结果
            print(f"Rerank服务调用异常：{str(e)}")
            reranked_docs = [{
                'doc_id': doc_id,
                'content': documents[doc_id],
                'score': score
            } for doc_id, score in candidate_docs[:FINAL_TOP_N]]
    else:
        reranked_docs = []
    
    # 获取最相关的文档
    retrieved_docs = reranked_docs[:FINAL_TOP_N]
    
    # 调用大模型生成答案
    context = '\n'.join([doc['content'] for doc in retrieved_docs])
    prompt = f"基于以下上下文回答问题：\n\n上下文：{context}\n\n问题：{query}\n\n回答："
    
    try:
        # 调用vllm提供的大模型服务
        response = requests.post(
            VLLM_SERVER_URL,
            json={'prompt': prompt, 'temperature': 0.7, 'max_tokens': 512}
        )
        
        if response.status_code == 200:
            answer = response.json().get('text', '')
        else:
            answer = '大模型服务调用失败'
            
    except Exception as e:
        answer = f'大模型服务调用异常：{str(e)}'
    
    return jsonify({
        'query': query,
        'answer': answer,
        'retrieved_docs': retrieved_docs
    })

@rag.route('/rag/docs', methods=['GET'])
def get_docs():
    """
    获取所有已索引的文档
    """
    return jsonify({
        'total_docs': len(documents),
        'documents': documents
    })

# 全局变量存储会议背景知识
meeting_backgrounds = {}

# 读取Word文档
@rag.route('/rag/upload_background', methods=['POST'])
def upload_background():
    """
    上传会议背景知识文件（Word/PPT/Excel）
    """
    global meeting_backgrounds
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file provided'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    # 生成唯一ID
    background_id = str(uuid.uuid4())
    
    # 使用统一存储系统保存文件
    temp_path = storage.save_file(file, directory='background', suffix=os.path.splitext(secure_filename(file.filename))[1].lower())
    
    # 解析文件内容
    content = ''
    file_ext = os.path.splitext(secure_filename(file.filename))[1].lower()
    
    try:
        if file_ext == '.docx':
            # 读取Word文档
            doc = Document(temp_path)
            for para in doc.paragraphs:
                content += para.text + '\n'
        elif file_ext == '.pptx':
            # 读取PPT文档
            prs = Presentation(temp_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        content += shape.text + '\n'
        elif file_ext == '.xlsx':
            # 读取Excel文档
            wb = load_workbook(temp_path)
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    if row and any(cell for cell in row):
                        row_content = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                        content += row_content + '\n'
        elif file_ext == '.txt':
            # 读取文本文件
            with open(temp_path, 'r', encoding='utf-8') as f:
                content = f.read()
        else:
            os.unlink(temp_path)
            return jsonify({'error': 'Unsupported file format. Supported: docx, pptx, xlsx, txt'}), 400
        
        # 保存背景知识
        meeting_backgrounds[background_id] = {
            'id': background_id,
            'filename': secure_filename(file.filename),
            'content': content,
            'upload_time': time.time()
        }
        
        os.unlink(temp_path)
        
        return jsonify({
            'message': 'Background knowledge uploaded and parsed successfully',
            'background_id': background_id,
            'filename': secure_filename(file.filename),
            'content_length': len(content)
        })
        
    except Exception as e:
        os.unlink(temp_path)
        return jsonify({'error': 'Failed to parse file: ' + str(e)}), 500

# 导入线程池和线程相关模块
from concurrent.futures import ThreadPoolExecutor
import threading





@rag.route('/rag/generate_minutes', methods=['POST'])
def generate_minutes():
    """
    生成会议记录：上传音频文件，结合背景知识生成会议记录
    """
    global meeting_backgrounds
    
    if 'audio' not in request.files:
        return jsonify({'error': 'No audio file provided'}), 400
    
    audio_file = request.files['audio']
    background_id = request.form.get('background_id')
    
    # 验证音频文件
    if audio_file.filename == '':
        return jsonify({'error': 'No audio file selected'}), 400
    
    # 检查背景知识
    background_content = ''
    if background_id:
        if background_id not in meeting_backgrounds:
            return jsonify({'error': 'Invalid background_id'}), 400
        background_content = meeting_backgrounds[background_id]['content']
    
    # 使用统一存储系统保存音频文件
    audio_ext = os.path.splitext(secure_filename(audio_file.filename))[1].lower()
    audio_temp_path = storage.save_file(audio_file, directory='audio', suffix=audio_ext)
    
    try:
        # 异步调用Celery任务生成会议记录
        task = generate_meeting_minutes.delay(audio_temp_path, audio_file.filename, background_id, background_content)
        
        return jsonify({
            'message': '会议记录生成任务已启动',
            'task_id': task.id,
            'background_id': background_id
        })
        
    except Exception as e:
        os.unlink(audio_temp_path)
        return jsonify({'error': 'Failed to start meeting minutes generation task: ' + str(e)}), 500

@rag.route('/rag/meeting_task/<task_id>', methods=['GET'])
def get_meeting_task_status(task_id):
    """
    查询会议记录生成任务状态
    """
    from tasks import generate_meeting_minutes
    
    # 使用Celery的AsyncResult查询任务状态
    task_result = generate_meeting_minutes.AsyncResult(task_id)
    
    if not task_result:
        return jsonify({'error': 'Invalid task_id'}), 400
    
    status = task_result.status.lower()
    message = ''
    transcript = ''
    minutes = ''
    background_id = ''
    
    # 获取任务结果
    if task_result.ready():
        if task_result.successful():
            result = task_result.result
            message = '会议记录生成成功'
            transcript = result.get('transcript', '')
            minutes = result.get('minutes', '')
            background_id = result.get('background_id', '')
        else:
            message = f'生成失败：{str(task_result.result)}'
    else:
        message = '任务正在处理中'
    
    return jsonify({
        'task_id': task_id,
        'status': status,
        'message': message,
        'transcript': transcript,
        'minutes': minutes,
        'background_id': background_id
    })